import io
import numpy as np

#
# NOTE: this imports the python-pptx module, not ourselves...
#
import pptx

import iwp.analysis
import iwp.labels
import iwp.rendering

# size of Matplotlib rendered XY slices, in inches.
XY_SLICE_WIDTH_INCHES  = 3.5
XY_SLICE_HEIGHT_INCHES = 3.9

# vertical offset, relative to top of the slide, of the rendered XY slices, in
# inches.
XY_SLICE_OFFSET_TOP_INCHES = 1.28

def _add_xy_slice_shape_group( slide, xy_slice_position, xy_slice_image, xy_slice_axes_position, variable_name, iwp_labels, y_axis_label_flag=False ):
    """
    Adds an XY slice image, its colorbar, and axes label decorations to an existing
    slide. Optionally overlays IWP labels onto the XY slice image.  All generated
    shapes are added to a new shape group for easy manipulation.

    NOTE: This function updates the slide object provided.

    Takes 7 arguments:

      slide                  - pptx.slide.Slide object to add XY slices to.  This is
                               modified during execution.
      xy_slice_position      - Tuple specifying the left, top, width, and height of the
                               XY slices to add to slide.  Must be an object derived
                               from pptx.util.Length.
      xy_slice_image         - XY slice image to add to the slide.  Specified either
                               as a path or a buffer of serialized bytes.
      xy_slice_axes_position - Tuple specifying the left, top, width, and height of
                               the XY slice image's underlying axes.  Must be an object
                               derived from pptx.util.Length.
      variable_name          - Name of the variable associated with xy_slice_image.
      iwp_labels             - List of IWP labels whose bounding boxes are normalized
                               into the range of [0, 1].  Each IWP label is rendered
                               as an unfilled box on the XY slice image.
      y_axis_label_flag      - Optional flag specifying whether the Y-axis label
                               should be generated.  Specify as True to the first, and
                               False to subsequent XY slices, if multiple are to be
                               added to a single slide to conserve horizontal space.
                               If omitted, defaults to False.

    Returns 1 value:

      xy_slice_group - The created pptx.shapes.group.GroupShape containing all of
                       the shapes associated with the XY slice image.

    """

    #
    # NOTE: all of the hard coded offsets and positions were copied or
    #       derived from manual layouts.
    #

    # name our coordinates so we can use them to derive other positions and
    # sizes.
    (xy_slice_left,
     xy_slice_top,
     xy_slice_width,
     xy_slice_height) = xy_slice_position

    xy_slice_group = slide.shapes.add_group_shape()

    # add the XY slice with a black border.
    xy_slice_picture = xy_slice_group.shapes.add_picture( xy_slice_image,
                                                          xy_slice_left,
                                                          xy_slice_top )

    xy_slice_picture.width  = xy_slice_width
    xy_slice_picture.height = xy_slice_height
    xy_slice_picture.name   = "XY Slice - {:s}".format(
        variable_name )

    # title the slice using the variable name.
    xy_slice_title_left   = xy_slice_left + pptx.util.Inches( 1.17 )
    xy_slice_title_top    = pptx.util.Inches( 1.31 )
    xy_slice_title_width  = pptx.util.Inches( 0.86 )
    xy_slice_title_height = pptx.util.Inches( 0.4 )

    xy_slice_title = xy_slice_group.shapes.add_textbox( xy_slice_title_left,
                                                        xy_slice_title_top,
                                                        xy_slice_title_width,
                                                        xy_slice_title_height )

    # bold, 16pt, and centered.
    p           = xy_slice_title.text_frame.paragraphs[0]
    p.text      = iwp.analysis.variable_name_to_title( variable_name,
                                                       latex_flag=False )
    p.font.bold = True
    p.font.size = pptx.util.Pt( 16 )
    p.alignment = pptx.enum.text.PP_ALIGN.CENTER

    # X-axis label.
    xy_slice_xaxis_left   = xy_slice_left + pptx.util.Inches( 1.17 )
    xy_slice_xaxis_top    = pptx.util.Inches( 4.94 )
    xy_slice_xaxis_width  = pptx.util.Inches( 0.86 )
    xy_slice_xaxis_height = pptx.util.Inches( 0.4 )

    xy_slice_xaxis = xy_slice_group.shapes.add_textbox( xy_slice_xaxis_left,
                                                        xy_slice_xaxis_top,
                                                        xy_slice_xaxis_width,
                                                        xy_slice_xaxis_height )
    xy_slice_xaxis.rotation = 0

    # bold, 14pt, and centered.
    p           = xy_slice_xaxis.text_frame.paragraphs[0]
    p.text      = "x/D"
    p.font.bold = True
    p.font.size = pptx.util.Pt( 14 )
    p.alignment = pptx.enum.text.PP_ALIGN.CENTER

    # Y-axis label.
    if y_axis_label_flag:
        #
        # NOTE: we don't position this relative to anything as it is assumed
        #       there will be at most one Y-axis label.
        #
        xy_slice_yaxis_left   = xy_slice_left - pptx.util.Inches( .31 )
        xy_slice_yaxis_top    = pptx.util.Inches( 3.03 )
        xy_slice_yaxis_width  = pptx.util.Inches( 0.86 )
        xy_slice_yaxis_height = pptx.util.Inches( 0.4 )

        xy_slice_yaxis = xy_slice_group.shapes.add_textbox( xy_slice_yaxis_left,
                                                            xy_slice_yaxis_top,
                                                            xy_slice_yaxis_width,
                                                            xy_slice_yaxis_height )
        xy_slice_yaxis.rotation = 270

        # bold, 14pt, and centered.
        p           = xy_slice_yaxis.text_frame.paragraphs[0]
        p.text      = "y/D"
        p.font.bold = True
        p.font.size = pptx.util.Pt( 14 )
        p.alignment = pptx.enum.text.PP_ALIGN.CENTER

    # flip the coordinate system for the labels.  they have the origin in the
    # bottom left, though the pptx module uses an origin in the upper left.  we
    # make a copy of the labels so we don't modify the caller's data.
    iwp_labels = iwp.labels.flipud_iwp_label_coordinates( iwp_labels,
                                                          1.0,
                                                          in_place_flag=False )

    # generate bounding boxes for each of the labels supplied.
    for iwp_label in iwp_labels:

        # convert the normalized label's four corners into (left, top, width,
        # height) so we can create a box representing the label.
        normalized_label_left   = iwp_label["bbox"]["x1"]
        normalized_label_top    = iwp_label["bbox"]["y1"]
        normalized_label_width  = (iwp_label["bbox"]["x2"] -
                                   iwp_label["bbox"]["x1"])
        normalized_label_height = (iwp_label["bbox"]["y2"] -
                                   iwp_label["bbox"]["y1"])

        # convert the normalized label's coordinates into the XY slice
        # picture's coordinates and position it.
        label_box_left = (xy_slice_left + xy_slice_axes_position[0] +
                          int( xy_slice_axes_position[2] *
                               normalized_label_left ))
        label_box_top  = (xy_slice_top + xy_slice_axes_position[1] +
                          int( xy_slice_axes_position[3] *
                               normalized_label_top ))

        # properly size the label by scaling it by the image dimensions.
        label_box_width  = int( xy_slice_axes_position[2] * normalized_label_width )
        label_box_height = int( xy_slice_axes_position[3] * normalized_label_height )

        # white border for an unfilled rectangle.  1 pt line thickness.
        label_box = xy_slice_group.shapes.add_shape( pptx.enum.shapes.MSO_SHAPE.RECTANGLE,
                                                     label_box_left,
                                                     label_box_top,
                                                     label_box_width,
                                                     label_box_height )

        label_box.fill.background()
        label_box.line.color.rgb = pptx.dml.color.RGBColor( 0xFF, 0xFF, 0xFF )
        label_box.line.width     = pptx.util.Pt( 1 )

        # IWP label name label.  add the label's name so that it is positioned
        # slightly above the label itself.
        #
        # NOTE: the offset between the edge of a text box and its first character is
        #       large enough that we need to position the label above and left of
        #       the label's upper-left corner.  without this, it is positioned
        #       too far to the right and looks wrong for small labels.
        #
        label_name_left   = label_box_left - pptx.util.Inches( 0.09 )
        label_name_top    = label_box_top  - pptx.util.Inches( 0.16 )
        label_name_width  = pptx.util.Inches( 0.25 )
        label_name_height = pptx.util.Inches( 0.25 )

        label_name = xy_slice_group.shapes.add_textbox( label_name_left,
                                                        label_name_top,
                                                        label_name_width,
                                                        label_name_height )
        label_name.rotation = 0

        # bold, 5pt, and left-aligned.
        p           = label_name.text_frame.paragraphs[0]
        p.text      = iwp.labels.get_iwp_label_name( iwp_label,
                                                     shortened_flag=True )
        p.font.color.rgb = pptx.dml.color.RGBColor( 0xFF, 0xFF, 0xFF )
        p.font.bold = True
        p.font.size = pptx.util.Pt( 5 )
        p.alignment = pptx.enum.text.PP_ALIGN.LEFT

    return xy_slice_group

def create_data_review_presentation( iwp_dataset, experiment_name, variable_names, time_xy_slice_pairs, data_limits, color_map, quantization_table_builder, iwp_labels=[] ):
    """
    Creates a Powerpoint presentation containing data review slides for a set of
    XY slices.  One slide per XY slice is generated, with up to three variables of
    the slice rendered side-by-side for comparison and review.  IWP labels may be
    overlaid on each XY slice rendered, and slices are normalized, quantized, and
    colored according to the caller's specification.

    Raises ValueError if too few or too many variables are specified.  This prevents
    the generated XY slices from being scaled down so much to be of no use.

    Takes 8 arguments:

      iwp_dataset                - iwp.data_loader.IWPDataset object containing the XY
                                   slices to generate review slides for.
      experiment_name            - String specifying the experiment that generated the slice.
      variable_names             - List of one, two, or three variable names to generate
                                   review images for.
      time_xy_slice_pairs        - List of tuples, (time index, XY slice index), specifying
                                   the XY slices to generate slides for.
      data_limits                - Dictionary of tuples, keys are variable names, keys are
                                   (minimum, maximum, standard deviation), specifying the
                                   global statistics to normalize each XY slice with.  If
                                   specified as None, statistics are computed for each XY
                                   slice independently.
      color_map                  - Matplotlib color map to apply to the underlying data.
      quantization_table_builder - Function that generates a quantization table when supplied
                                   four arguments: number of quantization levels, data minimum,
                                   data maximum, and data standard deviation.
      iwp_labels                 - Optional sequence of IWP labels to overlay on the
                                   generated data.

    Returns 1 value:

      presentation - pptx.Presentation object containing the slides generated.

    """

    # bail if we have too few or too many variables to review.  we have limited
    # space on each slide and cannot fit more than three XY slices before
    # running out of room.
    if not (0 < len( variable_names ) < 4):
        raise ValueError( "Must have between 1 and 3 variables to generate "
                          "review data for, received {:d}.".format(
                              len( variable_names ) ) )

    presentation = pptx.Presentation()

    # set a widescreen ratio (16:9) for this presentation's slides.
    presentation.slide_width  = pptx.util.Inches( 10 )
    presentation.slide_height = pptx.util.Inches( 5.625 )

    # start with a title-only slide.
    blank_slide_layout = presentation.slide_layouts[5]

    # specify the locations of each of the XY slices we're reviewing.  these
    # change as a function of slice count so that we use the most of the
    # available space and produce aesthetically pleasing results.  stored
    # as a list of position tuples (left, top, width, height), one per
    # XY slice to layout.
    if len( variable_names ) == 1:
        # single image centered in the middle of the slide.
        xy_slice_positions = [(pptx.util.Inches( 2.83 ),
                               pptx.util.Inches( XY_SLICE_OFFSET_TOP_INCHES ),
                               pptx.util.Inches( XY_SLICE_WIDTH_INCHES ),
                               pptx.util.Inches( XY_SLICE_HEIGHT_INCHES ))]
    elif len( variable_names ) == 2:
        # two images centered about the middle of the slide.
        xy_slice_positions = [(pptx.util.Inches( 1.0 ),
                               pptx.util.Inches( XY_SLICE_OFFSET_TOP_INCHES ),
                               pptx.util.Inches( XY_SLICE_WIDTH_INCHES ),
                               pptx.util.Inches( XY_SLICE_HEIGHT_INCHES )),
                              (pptx.util.Inches( 5.5 ),
                               pptx.util.Inches( XY_SLICE_OFFSET_TOP_INCHES ),
                               pptx.util.Inches( XY_SLICE_WIDTH_INCHES ),
                               pptx.util.Inches( XY_SLICE_HEIGHT_INCHES ))]
    elif len( variable_names ) == 3:
        # three images equally distributed across the slide.
        xy_slice_positions = [(pptx.util.Inches( 0 ),
                               pptx.util.Inches( XY_SLICE_OFFSET_TOP_INCHES ),
                               pptx.util.Inches( XY_SLICE_WIDTH_INCHES ),
                               pptx.util.Inches( XY_SLICE_HEIGHT_INCHES )),
                              (pptx.util.Inches( 3.25 ),
                               pptx.util.Inches( XY_SLICE_OFFSET_TOP_INCHES ),
                               pptx.util.Inches( XY_SLICE_WIDTH_INCHES ),
                               pptx.util.Inches( XY_SLICE_HEIGHT_INCHES )),
                              (pptx.util.Inches( 6.5 ),
                               pptx.util.Inches( XY_SLICE_OFFSET_TOP_INCHES ),
                               pptx.util.Inches( XY_SLICE_WIDTH_INCHES ),
                               pptx.util.Inches( XY_SLICE_HEIGHT_INCHES ))]

    # build a mapping from (time, xy slice) to labels so it is easy to identify
    # the relevant ones when building data review slides.
    iwp_labels_map = {}
    for iwp_label in iwp_labels:
        label_key = iwp.labels.get_iwp_label_key( iwp_label )

        # add this label to its (time, slice) list.
        if label_key not in iwp_labels_map:
            iwp_labels_map[label_key] = [iwp_label]
        else:
            iwp_labels_map[label_key].append( iwp_label )

    # iterate through each of the requested XY slices and make a slide for it.
    for time_index, xy_slice_index in time_xy_slice_pairs:
        current_slide = presentation.slides.add_slide( blank_slide_layout )

        # set the title.
        slide_title      = current_slide.placeholders[0]
        slide_title.text = "{:s}: Z={:03d}, Nt={:03d}".format(
            experiment_name,
            xy_slice_index,
            time_index )

        # pull the data for this XY slice.
        xy_slice_array = iwp_dataset.get_xy_slice( time_index,
                                                   xy_slice_index )

        # construct a label key so we can lookup the labels associated with
        # this XY slice.
        label_key = (time_index, xy_slice_index)

        # iterate through each of the variables and add a group containing
        # the rendered data with titles and axes labels.  size and positions
        # are provided allowing for variable count-specific layouts (e.g.
        # centered, big images for a single variable vs smaller, multi-column
        # layouts for multiple variables).
        for variable_index, variable_name in enumerate( variable_names ):
            # get this variable's statistics so we can quantize it properly.
            if data_limits is not None:
                # pull our variable's statistics out of the global statistics.
                variable_statistics = (data_limits[variable_name]["minimum"],
                                       data_limits[variable_name]["maximum"],
                                       data_limits[variable_name]["standard_deviation"])
            else:
                # compute our variable's local statistics from the current XY
                # slice.
                variable_statistics = (xy_slice_array[variable_index, :].min(),
                                       xy_slice_array[variable_index, :].max(),
                                       xy_slice_array[variable_index, :].std())

            quantization_table = quantization_table_builder( 256,
                                                             *variable_statistics )

            # compute a scale factor (order of magnitude) for the colorbar
            # ticks.  all tick labels are of the magnitude computed here.
            #
            # NOTE: we choose the floor of the extrema's log10 magnitude to
            #       favor extremal tick values that are whole number with a
            #       smaller exponent, rather than a decimal number with a larger
            #       exponent.  (e.g. +-5 x 10^-2 vs +-.5 x 10^-3).
            #
            oom_factor = np.floor( np.log10( np.max( np.abs( variable_statistics[:2] ) ) ) )

            # format our colorbar tick labels to always have our scale factor displayed
            # and render as "x 10^<exponent>" instead of "1e<exponent>".
            colorbar_formatter = iwp.analysis.FixedScientificFormatter( oom_factor,
                                                                        "%1.1f",
                                                                        offset_flag=True,
                                                                        math_text_flag=True )

            # render this XY slice to an image.  we use Matplotlib so we get
            # properly labeled axes and a colorbar, as well as consistency with
            # other IWP visualization workflows.
            #
            # we hang onto the underlying figure handle so we can access the
            # underlying axes to identify where the XY slice data is relative
            # to the rendered image.  this enables us to overlay our labels
            # in the correct location.
            (xy_slice_image,
             xy_slice_fig_h) = iwp.rendering.array_to_image_imshow( xy_slice_array[variable_index, :],
                                                                    quantization_table,
                                                                    color_map,
                                                                    figure_size=(xy_slice_positions[variable_index][2].inches,
                                                                                 xy_slice_positions[variable_index][3].inches),
                                                                    show_axes_labels_flag=False,
                                                                    colorbar_flag=True,
                                                                    colorbar_formatter=colorbar_formatter,
                                                                    constrained_layout_flag=False )

            # set the white background to transparent to work around the fact
            # that we have wide images that overlap.
            #
            # NOTE: this is done to maximize the horizontal real estate on the
            #       slide but sometime results in a colorbar tick label being
            #       hidden by the excess horizontal margin on the image to its
            #       right (in the case of two and three image layouts).
            #
            xy_slice_image = iwp.rendering.image_make_white_transparent( xy_slice_image )

            xy_slice_image_buffer = io.BytesIO()
            xy_slice_image.save( xy_slice_image_buffer, format="png" )

            # get the figure's size so we can properly scale it and position
            # pptx labels onto it.
            xy_slice_figure_size = xy_slice_fig_h.get_size_inches()

            # get the XY slice axes and the bounding box of the rendered data,
            # relative to its parent figure.
            #
            # NOTE: the XY slice is the first axes in the figure.  its colorbar
            #       is second.
            #
            # NOTE: axes positions are normalized figure coordinates with an
            #       origin in the bottom left (!) like so:
            #
            #
            #       axes y=1 ----------------------------------------
            #
            #                 (y1, x0)                  (y1, x1)
            #                     +                         +
            #
            #                     +                         +
            #                 (y0, x0)                  (y0, x0)
            #
            #       axes y=0 ----------------------------------------
            #
            #       one wants (1-y1, x0) as the offset from the top-left corner
            #       of the XY slice figure to the top-left corner of the XY
            #       slice axes.
            #
            xy_slice_ax_h      = xy_slice_fig_h.get_axes()[0]
            xy_slice_axes_bbox = xy_slice_ax_h.get_position()

            # compute offsets within the rendered figure to the XY slice data
            # itself.  we need this so we can correctly position label boxes
            # since the data don't start at the corner of the figure.
            #
            # NOTE: mind the implicit flip up/down that accounts for the
            #       axes coordinate system not matching pptx's.
            #
            xy_slice_axes_offset_x = pptx.util.Inches( xy_slice_axes_bbox.x0 *
                                                       xy_slice_figure_size[0] )
            xy_slice_axes_offset_y = pptx.util.Inches( (1 - xy_slice_axes_bbox.y1) *
                                                       xy_slice_figure_size[1] )

            # compute the size of the XY slice data within the rendered figure.
            xy_slice_axes_width  = pptx.util.Inches( xy_slice_axes_bbox.width *
                                                     xy_slice_figure_size[0] )
            xy_slice_axes_height = pptx.util.Inches( xy_slice_axes_bbox.height *
                                                     xy_slice_figure_size[1] )

            # add this XY slice to the slide in a group.  only generate the
            # y-axis labeling on the first image so we efficiently use our
            # horizontal space and avoid clutter.
            _add_xy_slice_shape_group( current_slide,
                                       xy_slice_positions[variable_index],
                                       xy_slice_image_buffer,
                                       (xy_slice_axes_offset_x,
                                        xy_slice_axes_offset_y,
                                        xy_slice_axes_width,
                                        xy_slice_axes_height),
                                       variable_name,
                                       iwp_labels_map.get( label_key, [] ),
                                       y_axis_label_flag=(variable_index == 0))

    return presentation
